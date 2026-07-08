"""
Generate PULSE pitch deck v2 — 12-slide framework edition (.pptx).
Run: python scripts/generate_pulse_pitch_deck_v2.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Global design standards
PRIMARY = RGBColor(0x0F, 0x4C, 0x81)
TEAL = RGBColor(0x00, 0xA6, 0xA6)
ACCENT = RGBColor(0x16, 0xA3, 0x4A)
WARNING = RGBColor(0xF5, 0x9E, 0x0B)
CRITICAL = RGBColor(0xDC, 0x26, 0x26)
CANVAS = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BLUE = RGBColor(0xE0, 0xF2, 0xFE)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 12
OUT_PATH = Path(__file__).resolve().parents[1] / "PULSE_Pitch_Deck_Framework.pptx"

# Font fallbacks (Plus Jakarta Sans / Inter may not be installed; Calibri is safe fallback)
HEADLINE_FONT = "Calibri"
BODY_FONT = "Calibri"


def set_bg(slide, color=CANVAS):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, size=12, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    p.alignment = align
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = HEADLINE_FONT if bold and size >= 14 else BODY_FONT
    return box, tf


def bullet(tf, text, size=11, bold=False, color=DARK, space=4, level=0):
    p = tf.add_paragraph()
    p.level = level
    p.space_before = Pt(space)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = BODY_FONT
    return p


def notes(slide, script):
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    tf.text = script


def header_bar(slide):
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.1), PRIMARY)


def label(slide, text):
    rect(slide, Inches(0.5), Inches(0.28), Inches(0.07), Inches(0.4), TEAL)
    txt(slide, Inches(0.65), Inches(0.22), Inches(5), Inches(0.35),
        text.upper(), size=10, bold=True, color=TEAL)


def footer(slide, n):
    txt(slide, Inches(0.5), Inches(7.05), Inches(11.5), Inches(0.28),
        "PULSE · Bago City College · BSIS · Framework Pitch 2026",
        size=8, color=MUTED, align=PP_ALIGN.CENTER)
    txt(slide, Inches(12.2), Inches(7.05), Inches(0.8), Inches(0.28),
        f"{n}/{TOTAL}", size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def title_block(slide, title, subtitle=None):
    txt(slide, Inches(0.55), Inches(0.75), Inches(12), Inches(0.7),
        title, size=30, bold=True, color=PRIMARY)
    if subtitle:
        txt(slide, Inches(0.55), Inches(1.35), Inches(12), Inches(0.4),
            subtitle, size=13, color=MUTED)


def slide_01(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide)
    rect(slide, Inches(7.8), Inches(0.1), Inches(5.53), Inches(7.4), PRIMARY)
    rect(slide, Inches(8.1), Inches(1.0), Inches(4.9), Inches(5.5), RGBColor(0x0A, 0x3A, 0x62))

    txt(slide, Inches(0.6), Inches(1.0), Inches(7), Inches(0.9),
        "PULSE", size=56, bold=True, color=PRIMARY)
    txt(slide, Inches(0.6), Inches(1.85), Inches(7.1), Inches(1.0),
        "Real-Time Syndromic Surveillance and Risk Forecasting using "
        "Isolation Forest Algorithms and Geospatial Data Fusion",
        size=15, bold=True, color=DARK)
    txt(slide, Inches(0.6), Inches(2.95), Inches(7.1), Inches(0.5),
        "Intelligent Epidemiological Surveillance for a Resilient Future",
        size=14, bold=True, color=TEAL, italic=True)

    rect(slide, Inches(0.6), Inches(3.65), Inches(7), Inches(0.02), TEAL)
    txt(slide, Inches(0.6), Inches(3.85), Inches(7), Inches(0.35),
        "Target Domain", size=10, bold=True, color=MUTED)
    txt(slide, Inches(0.6), Inches(4.15), Inches(7), Inches(0.45),
        "Bago City Health Office — Community-Based Disease Surveillance Unit",
        size=13, bold=True, color=PRIMARY)
    txt(slide, Inches(0.6), Inches(4.7), Inches(7), Inches(0.35),
        "Academic Anchor", size=10, bold=True, color=MUTED)
    txt(slide, Inches(0.6), Inches(5.0), Inches(7), Inches(0.4),
        "Bago City College — College of Information Systems",
        size=13, bold=True, color=DARK)
    txt(slide, Inches(0.6), Inches(5.65), Inches(7), Inches(0.4),
        "Banez · Betita · Sarcia", size=11, bold=True, color=MUTED)

    txt(slide, Inches(8.3), Inches(1.3), Inches(4.5), Inches(4.8),
        "HERO VISUAL\n\n24-Barangay\nAPTAS Choropleth Map\n\n"
        "● Critical / High zones\n● Case pins & heatmap\n● Real-time risk tiers",
        size=12, color=WHITE, align=PP_ALIGN.CENTER)

    notes(slide,
          "Judges, good day. We are the team behind PULSE, presenting an intelligent "
          "epidemiological surveillance system built for a resilient future. Our project "
          "re-engineers how local government units detect and handle emerging health crises. "
          "Developed in cooperation with the Bago City Health Office, PULSE bridges the gap "
          "between field-level community symptoms and automated, city-wide protective actions.")


def slide_02(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide)
    label(slide, "The Problem")
    title_block(slide, "The Problem",
                "Disaster Risk Reduction and Public Health Officers cannot deliver timely services because:")

    problems = [
        ("Inefficient Registration and Record Management",
         "Disease surveillance relies on manual, paper-based CDSS worksheets, leading to slow data consolidation and human error.",
         CRITICAL),
        ("Lack of Reliable Tracking and Information Systems",
         "Current setups only present raw numerical case counts without automated trend metrics or spatial risk analytics.",
         WARNING),
        ("Limited Dissemination of Warnings and Options",
         "The complete absence of an automated alert system delays early outbreak identification and restricts proactive health responses.",
         PRIMARY),
    ]
    y = 1.85
    for i, (title, desc, color) in enumerate(problems, 1):
        rect(slide, Inches(0.55), Inches(y), Inches(12.2), Inches(1.35), WHITE, BORDER)
        rect(slide, Inches(0.55), Inches(y), Inches(0.14), Inches(1.35), color)
        txt(slide, Inches(0.85), Inches(y + 0.12), Inches(0.5), Inches(0.35),
            f"0{i}", size=16, bold=True, color=color)
        txt(slide, Inches(1.35), Inches(y + 0.1), Inches(11), Inches(0.4),
            title, size=13, bold=True, color=DARK)
        txt(slide, Inches(1.35), Inches(y + 0.52), Inches(11), Inches(0.7),
            desc, size=11, color=MUTED)
        y += 1.5

    footer(slide, 2)
    notes(slide,
          "Every health crisis starts with small, localized signals in our communities, but public health "
          "officers are unable to deliver timely interventions because of three critical bottlenecks. First, "
          "record management relies on slow, paper-based worksheets that cause massive data consolidation lags. "
          "Second, local units lack reliable tracking systems, viewing raw numbers with zero spatial visualization. "
          "Third, there is no automated early warning network, meaning outbreak patterns are caught only after "
          "they have already spread.")


def slide_03(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide)
    label(slide, "The Solution")
    title_block(slide, "PULSE Solution",
                "AI-Powered Public Health Ingestion and Surveillance System with Predictive Analytics")

    box = rect(slide, Inches(0.55), Inches(1.75), Inches(12.2), Inches(1.0), LIGHT_BLUE)
    txt(slide, Inches(0.75), Inches(1.9), Inches(11.8), Inches(0.8),
        "Core Value Proposition: Transforms delayed, paper-bound data gathering into real-time digital intelligence.",
        size=14, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

    points = [
        "Replaces manual, reactive data processing with continuous machine learning screening.",
        "Empowers local health administrations to transition from counting historical infections to executing proactive, preventative community containment strategies.",
        "Unifies syndromic intake, anomaly detection, geospatial fusion, and APTAS alerting in one platform.",
    ]
    y = 3.0
    for p in points:
        rect(slide, Inches(0.55), Inches(y), Inches(12.2), Inches(0.85), WHITE, BORDER)
        txt(slide, Inches(0.85), Inches(y + 0.2), Inches(11.6), Inches(0.55),
            f"✓  {p}", size=12, color=DARK)
        y += 1.0

    rect(slide, Inches(0.55), Inches(6.1), Inches(12.2), Inches(0.65), PRIMARY)
    txt(slide, Inches(0.75), Inches(6.25), Inches(11.8), Inches(0.4),
        "From reactive counting → proactive, preventative community containment",
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    footer(slide, 3)
    notes(slide,
          "Our answer is PULSE: an AI-powered public health ingestion and surveillance system featuring "
          "advanced machine learning and predictive analytics. PULSE eliminates the traditional paper trail. "
          "It gives cities an intelligent digital layer that continually scans field data, visually maps community "
          "risk factors, and shifts local health governance from a reactive, descriptive model to a proactive, "
          "preventative shield.")


def slide_04(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide)
    label(slide, "Product Capabilities")
    title_block(slide, "Product Capabilities", "Three automated modules that deliver end-to-end surveillance")

    modules = [
        ("Automated Registration &\nInformation Management",
         "Replaces paper forms with batch digital collection layouts for syndromic health data, securely archiving demographics and symptom strings into cloud relational storage.",
         TEAL),
        ("Real-Time Tracking &\nMonitoring System",
         "Runs background anomaly evaluations via the Isolation Forest Algorithm to calculate precise risk score values and flag abnormal symptom spikes instantly at intake.",
         PRIMARY),
        ("Digital Outbreak Route &\nHotspot Allocation System",
         "Leverages geospatial data fusion to map coordinates across 24 official boundaries, plotting dynamic heatmaps and triggering the APTAS early warning notification module.",
         ACCENT),
    ]
    x = 0.45
    for name, desc, color in modules:
        rect(slide, Inches(x), Inches(1.75), Inches(4.05), Inches(4.85), WHITE, BORDER)
        rect(slide, Inches(x), Inches(1.75), Inches(4.05), Inches(0.75), color)
        txt(slide, Inches(x + 0.15), Inches(1.82), Inches(3.75), Inches(0.7),
            name, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, Inches(x + 0.2), Inches(2.7), Inches(3.65), Inches(3.6),
            desc, size=11, color=DARK)
        x += 4.2

    footer(slide, 4)
    notes(slide,
          "We achieved this breakthrough by deploying three automated product modules. First, our Automated "
          "Information Management System digitizes the front line, giving field workers secure batch entry "
          "layouts for symptoms. Second, our Real-Time Tracking Core runs background Python algorithms utilizing "
          "the Isolation Forest model to flag abnormal anomalies the moment they hit the server. Third, our "
          "Geospatial Hotspot Allocation mapping overlays these alerts instantly onto interactive boundary layers, "
          "dispatching automated threshold warnings via the APTAS engine to city administrators.")


def slide_05_market(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide)
    label(slide, "Market Size")
    title_block(slide, "Market Size", "Philippines LGU health surveillance software opportunity")

    markets = [
        ("TAM", "₱4.7 Billion", "Philippines market space across ~1,600 municipal and city LGUs requiring modern health and disaster software infrastructure.", PRIMARY),
        ("SAM", "₱126 Million", "Negros Island regional scope covering approximately 50 active localized government units.", TEAL),
        ("SOM", "Bago City Launch", "Initial deployment and model optimization target localized strictly within Bago City's 24 health-reporting barangay divisions.", ACCENT),
    ]
    y = 1.85
    for tag, value, desc, color in markets:
        rect(slide, Inches(0.55), Inches(y), Inches(12.2), Inches(1.55), WHITE, BORDER)
        rect(slide, Inches(0.55), Inches(y), Inches(1.4), Inches(1.55), color)
        txt(slide, Inches(0.55), Inches(y + 0.55), Inches(1.4), Inches(0.4),
            tag, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, Inches(2.15), Inches(y + 0.2), Inches(3.5), Inches(0.5),
            value, size=22, bold=True, color=color)
        txt(slide, Inches(2.15), Inches(y + 0.75), Inches(10.3), Inches(0.65),
            desc, size=11, color=MUTED)
        y += 1.7

    footer(slide, 5)
    notes(slide,
          "This framework tackles a significant software space. Across the Philippines, local government "
          "infrastructure represents a 4.7-billion-peso market across 1,600 municipal units. Regionally, on "
          "Negros Island alone, that potential represents 126 million pesos across 50 LGUs. PULSE addresses "
          "this market directly, using Bago City's 24 reporting barangay divisions as our optimized launchpad "
          "to validate our software's commercial scaling capabilities.")


def slide_06_business(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide)
    label(slide, "Business Model")
    title_block(slide, "Business Model", "Core Framework: License-Based Revenue Model")

    items = [
        ("Perpetual Licensing", "One-time software installation package tailored for local government operations and centralized health units."),
        ("Implementation & Customization Fees", "Structural setup and matching localized GeoJSON maps to specific municipal administrative zones."),
        ("Maintenance & Support Contracts", "Yearly service agreements to oversee cloud database uptimes and guarantee high availability parameters."),
        ("Training & Capacity Building", "Professional on-site orientation and onboarding programs for field workers, encoders, and municipal surveillance officers."),
    ]
    y = 1.85
    for title, desc in items:
        rect(slide, Inches(0.55), Inches(y), Inches(12.2), Inches(1.15), WHITE, BORDER)
        rect(slide, Inches(0.55), Inches(y), Inches(0.1), Inches(1.15), TEAL)
        txt(slide, Inches(0.85), Inches(y + 0.15), Inches(11.5), Inches(0.35),
            title, size=13, bold=True, color=PRIMARY)
        txt(slide, Inches(0.85), Inches(y + 0.52), Inches(11.5), Inches(0.55),
            desc, size=11, color=MUTED)
        y += 1.28

    footer(slide, 6)
    notes(slide,
          "To guarantee operational longevity, we utilize a sustainable, license-based revenue model tailored "
          "for institutional procurement. We generate revenue through perpetual licensing implementation, paired "
          "with customization fees for mapping local municipal boundaries. This is backed by predictable, recurring "
          "maintenance support contracts and hands-on capacity-building programs to onboard field health workers seamlessly.")


def slide_07_gtm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide)
    label(slide, "Go-To-Market")
    title_block(slide, "Go-To-Market Strategy", "Three-phase rollout from Bago City to nationwide scale")

    phases = [
        ("Phase 1", "2026", "Pilot Testing",
         "Execute localized trial rollouts alongside high-risk barangays in Bago City to gather user insights and demonstrate functional real-time accuracy.",
         TEAL),
        ("Phase 2", "2027", "Regional Expansion",
         "Expand operations to provincial disaster and regional public health surveillance agencies across Negros Occidental via strategic partnerships.",
         PRIMARY),
        ("Phase 3", "2028", "Nationwide Scaling",
         "Scale the centralized web system nationwide via modular cloud architectures and localized dialect configurations for extensive municipal accessibility.",
         ACCENT),
    ]
    x = 0.55
    for phase, year, name, desc, color in phases:
        rect(slide, Inches(x), Inches(1.85), Inches(3.95), Inches(4.6), WHITE, BORDER)
        rect(slide, Inches(x), Inches(1.85), Inches(3.95), Inches(0.55), color)
        txt(slide, Inches(x + 0.15), Inches(1.92), Inches(3.65), Inches(0.4),
            f"{phase} · {year}", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, Inches(x + 0.15), Inches(2.55), Inches(3.65), Inches(0.45),
            name, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(slide, Inches(x + 0.2), Inches(3.1), Inches(3.55), Inches(3.1),
            desc, size=11, color=DARK)
        if x < 8:
            txt(slide, Inches(x + 4.0), Inches(3.8), Inches(0.4), Inches(0.4),
                "→", size=24, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
        x += 4.15

    footer(slide, 7)
    notes(slide,
          "Our strategic timeline is structured across three clear phases. We are starting right now in 2026 "
          "with localized pilot testing across high-risk barangays in Bago City to capture real-world feedback. "
          "In 2027, Phase 2 scales our platform to regional provincial disaster and health networks via administrative "
          "partnerships. By 2028, we will achieve nationwide scaling using modular cloud software architectures built "
          "to service municipal infrastructures throughout the country.")


def slide_08_competition(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide)
    label(slide, "Competition Matrix")
    title_block(slide, "Competition Matrix", "PULSE vs. existing surveillance approaches")

    # Table header
    cols = [Inches(2.4), Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.0)]
    headers = ["System", "Alerts", "Maps", "ML / Predictive", "Real-Time Ingestion"]
    x = 0.55
    for i, h in enumerate(headers):
        w = cols[i] if i < len(cols) else Inches(2.2)
        rect(slide, x, Inches(1.75), w, Inches(0.45), PRIMARY)
        txt(slide, x + Inches(0.05), Inches(1.8), w - Inches(0.1), Inches(0.35),
            h, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += w

    rows = [
        ("PULSE", "✓ Full APTAS", "✓ Choropleth + pins", "✓ Isolation Forest", "✓ Batch digital intake", ACCENT),
        ("PIDSR (Manual)", "✗ None", "✗ None", "✗ None", "✗ Paper → manual encode", MUTED),
        ("DHIS2", "✗ Limited", "△ Basic dashboards", "✗ None", "△ Manual entry", MUTED),
        ("Generic Web Maps", "✗ None", "△ Static pins", "✗ None", "✗ No health logic", MUTED),
    ]
    y = 2.25
    for name, a, m, ml, ing, highlight in rows:
        bg = RGBColor(0xEC, 0xFD, 0xF5) if highlight == ACCENT else WHITE
        vals = [name, a, m, ml, ing]
        x = 0.55
        for i, val in enumerate(vals):
            w = cols[i]
            rect(slide, x, Inches(y), w, Inches(0.95), bg, BORDER)
            txt(slide, x + Inches(0.08), Inches(y + 0.2), w - Inches(0.16), Inches(0.6),
                val, size=9, bold=(i == 0), color=DARK if highlight != ACCENT or i > 0 else PRIMARY)
            x += w
        y += 1.0

    # Detail bullets below table
    competitors = [
        ("PIDSR (Manual Framework)", "Validates national reporting standards, but features zero automated predictive capability or real-time spatial visualization."),
        ("DHIS2 Platform", "Offers solid data aggregation dashboards, but completely lacks AI-driven anomaly detection models or early warning alerts."),
        ("Generic Web Maps", "Provides simple, static pinpoint visualizations, but features no public health logic, threshold evaluation, or role-scoped data forms."),
    ]
    y = 6.35
    for title, desc in competitors[:1]:
        txt(slide, Inches(0.55), Inches(y), Inches(12), Inches(0.3),
            f"{title}: {desc}", size=8, color=MUTED)

    footer(slide, 8)
    notes(slide,
          "When looking at the competitive landscape, current baseline systems fail to deliver integrated, "
          "proactive capabilities. Traditional frameworks like manual PIDSR steps meet national documentation "
          "requirements, but rely entirely on slow, paper-heavy processing. Platforms like DHIS2 aggregate data well, "
          "but completely lack machine learning anomaly models or early warning alerts. PULSE is the only comprehensive "
          "solution that delivers real-time ingestion, automated anomaly tracking, predictive risk scoring, and spatial "
          "data fusion under one unified digital dashboard.")


def slide_09_ask(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide)
    label(slide, "The Ask")
    title_block(slide, "The Ask & Partnership Opportunity")

    rect(slide, Inches(0.55), Inches(1.75), Inches(5.8), Inches(2.2), PRIMARY)
    txt(slide, Inches(0.75), Inches(1.9), Inches(5.4), Inches(0.35),
        "Funding Ask", size=12, bold=True, color=TEAL)
    txt(slide, Inches(0.75), Inches(2.3), Inches(5.4), Inches(0.7),
        "₱4.2 Million", size=36, bold=True, color=WHITE)
    txt(slide, Inches(0.75), Inches(3.05), Inches(5.4), Inches(0.8),
        "Allocated for on-premise cloud server infrastructure setup, secure data clustering development, "
        "and perpetual licensing optimization by LGUs and NGOs.",
        size=11, color=WHITE)

    rect(slide, Inches(6.6), Inches(1.75), Inches(6.15), Inches(2.2), WHITE, BORDER)
    txt(slide, Inches(6.8), Inches(1.9), Inches(5.75), Inches(0.35),
        "Partnership Alignments", size=12, bold=True, color=PRIMARY)
    partners = [
        "• Local Government Units (LGUs)",
        "• Department of Health regional offices",
        "• Non-governmental organizations (NGOs)",
        "• Community emergency management operations",
    ]
    y = 2.35
    for p in partners:
        txt(slide, Inches(6.8), Inches(y), Inches(5.75), Inches(0.3), p, size=11, color=DARK)
        y += 0.38

    rect(slide, Inches(0.55), Inches(4.2), Inches(12.2), Inches(1.5), LIGHT_BLUE)
    txt(slide, Inches(0.75), Inches(4.35), Inches(11.8), Inches(0.35),
        "Strategic Mentorship Focus", size=12, bold=True, color=PRIMARY)
    txt(slide, Inches(0.75), Inches(4.75), Inches(11.8), Inches(0.8),
        "Refining automated machine learning model parameters and structural civic health resilience "
        "deployment strategies.",
        size=12, color=DARK)

    footer(slide, 9)
    notes(slide,
          "Today, we are presenting a partnership opportunity. We are seeking a funding ask of 4.2 million pesos "
          "to establish resilient on-premise and cloud server infrastructures, ensuring data security and high platform "
          "availability for local government units. We are looking to align with forward-thinking LGUs, NGOs, and "
          "healthcare networks, while welcoming expert technical mentorship to further calibrate our AI model scaling "
          "and maximize our long-term public health impact.")


def slide_10_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide)
    label(slide, "Demonstration")
    title_block(slide, "Product Demonstration Hub",
                "Fully functional web system — ISO/IEC 25010 Quality Model verified")

    rect(slide, Inches(0.55), Inches(1.75), Inches(12.2), Inches(0.55), RGBColor(0xEC, 0xFD, 0xF5))
    txt(slide, Inches(0.75), Inches(1.88), Inches(11.8), Inches(0.35),
        "✓  Verified against ISO/IEC 25010 — functional suitability, reliability, and data security",
        size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    demos = [
        ("Administrative Command Dashboard",
         "Displays global morbidity trends, active risk tiers, and real-time alert logs.",
         PRIMARY),
        ("Geospatial Outbreak Choropleth",
         "Pins cases dynamically, highlights hotspot groupings, and locks data views securely.",
         CRITICAL),
        ("System Facility Management",
         "Accessible, role-scoped interfaces built for Encoders, Health Workers, and Supervising Officers.",
         TEAL),
    ]
    x = 0.55
    for title, desc, color in demos:
        rect(slide, Inches(x), Inches(2.5), Inches(3.95), Inches(3.8), WHITE, BORDER)
        rect(slide, Inches(x), Inches(2.5), Inches(3.95), Inches(0.5), color)
        txt(slide, Inches(x + 0.15), Inches(2.58), Inches(3.65), Inches(0.4),
            title, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(slide, Inches(x + 0.2), Inches(3.15), Inches(3.55), Inches(2.0),
             RGBColor(0xF1, 0xF5, 0xF9))
        txt(slide, Inches(x + 0.2), Inches(3.8), Inches(3.55), Inches(0.5),
            "[ Screenshot ]", size=9, color=MUTED, align=PP_ALIGN.CENTER)
        txt(slide, Inches(x + 0.2), Inches(5.3), Inches(3.55), Inches(0.85),
            desc, size=10, color=DARK)
        x += 4.15

    footer(slide, 10)
    notes(slide,
          "This project is not a concept—it is a functional web system tested directly against international "
          "ISO/IEC 25010 engineering standards to guarantee performance security and operational reliability. "
          "Our administrative platform unifies live morbidity counts with predictive analytics. Our interactive "
          "maps instantly group cluster locations into actionable threat tiers. From secure login screens to localized "
          "fields for frontline workers, the platform is designed to convert raw symptoms into protective institutional knowledge.")


def slide_11_impact(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide)
    label(slide, "Impact & Relevance")
    title_block(slide, "Impact & Relevance", "Beyond technology — lasting public health outcomes")

    impacts = [
        ("Proactive Public Health Management",
         "PULSE enables instantaneous, data-driven decision-making, changing community surveillance from delayed to predictive.",
         PRIMARY),
        ("Efficient Surveillance Procedures",
         "The platform automates data ingestion and analysis, completely eliminating old manual filing bottlenecks.",
         TEAL),
        ("Community Resilience Building",
         "Drastically improves cross-agency response coordination, fostering deep civic trust and protecting local populations.",
         ACCENT),
        ("Social Impact & Sustainability",
         "Supports long-term safety strategies and structural outbreak preparedness frameworks that extend well beyond simple technology.",
         WARNING),
    ]
    y = 1.85
    for title, desc, color in impacts:
        rect(slide, Inches(0.55), Inches(y), Inches(12.2), Inches(1.15), WHITE, BORDER)
        rect(slide, Inches(0.55), Inches(y), Inches(0.12), Inches(1.15), color)
        txt(slide, Inches(0.85), Inches(y + 0.15), Inches(11.5), Inches(0.35),
            title, size=13, bold=True, color=DARK)
        txt(slide, Inches(0.85), Inches(y + 0.52), Inches(11.5), Inches(0.55),
            desc, size=11, color=MUTED)
        y += 1.28

    footer(slide, 11)
    notes(slide,
          "The overall impact of PULSE stretches far beyond code architecture. We introduce proactive health "
          "management, allowing cities to neutralize health risks before they spread. We optimize institutional "
          "efficiency by removing manual data entry bottlenecks. Ultimately, we build deep community resilience. "
          "By providing transparent, accurate, and rapid responses, PULSE helps protect vulnerable neighborhoods and "
          "ensures long-term social safety for our cities.")


def slide_12_team(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    header_bar(slide)
    label(slide, "The Team")
    title_block(slide, "Founding Team & Closing")

    team = [
        ("Hernani J. Banez", "Project Manager & Backend Infrastructure Specialist"),
        ("Jaykirt M. Betita", "Lead Programmer & Core Machine Learning Analytics Engineer"),
        ("Earl Adrian M. Sarcia", "UI/UX Designer & Geospatial Data Analyst"),
    ]
    x = 0.55
    for name, role in team:
        rect(slide, Inches(x), Inches(1.75), Inches(3.95), Inches(2.3), WHITE, BORDER)
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.45), Inches(1.95), Inches(1.0), Inches(1.0))
        oval.fill.solid()
        oval.fill.fore_color.rgb = PRIMARY
        oval.line.fill.background()
        txt(slide, Inches(x + 0.15), Inches(3.1), Inches(3.65), Inches(0.35),
            name, size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        txt(slide, Inches(x + 0.15), Inches(3.45), Inches(3.65), Inches(0.5),
            role, size=9, color=MUTED, align=PP_ALIGN.CENTER)
        x += 4.15

    rect(slide, Inches(0.55), Inches(4.25), Inches(12.2), Inches(1.1), LIGHT_BLUE)
    txt(slide, Inches(0.75), Inches(4.4), Inches(11.8), Inches(0.85),
        "Institutional Alignment: Bago City College — BS Information Systems Department\n"
        "Project Mentor & Coordinator: Anthony S. Malabanan, MIT, MAT-Math",
        size=11, color=PRIMARY)

    banner = rect(slide, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.2), PRIMARY)
    txt(slide, Inches(0.75), Inches(5.75), Inches(11.8), Inches(0.85),
        "Intelligent Disease Surveillance and Risk Management for a Resilient Future.",
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(slide, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.35),
        "Thank you. We welcome your questions.", size=12, bold=True,
        color=TEAL, align=PP_ALIGN.CENTER)

    footer(slide, 12)
    notes(slide,
          "We are Banez, Betita, and Sarcia from Bago City College. PULSE closes with our vision: "
          "Intelligent Disease Surveillance and Risk Management for a Resilient Future. Thank you.")


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05_market(prs)
    slide_06_business(prs)
    slide_07_gtm(prs)
    slide_08_competition(prs)
    slide_09_ask(prs)
    slide_10_demo(prs)
    slide_11_impact(prs)
    slide_12_team(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
