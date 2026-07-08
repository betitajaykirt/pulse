"""
Generate PULSE Innovation Hackathon pitch deck (.pptx).
Run: python scripts/generate_pulse_pitch_deck.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# PULSE design tokens
PRIMARY = RGBColor(0x0F, 0x4C, 0x81)
TEAL = RGBColor(0x00, 0xA6, 0xA6)
ACCENT = RGBColor(0x16, 0xA3, 0x4A)
WARNING = RGBColor(0xF5, 0x9E, 0x0B)
CRITICAL = RGBColor(0xDC, 0x26, 0x26)
CANVAS = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
OUT_PATH = Path(__file__).resolve().parents[1] / "PULSE_Pitch_Deck.pptx"


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text="", size=14, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return box, tf


def add_paragraph(tf, text, size=14, bold=False, color=DARK, space_before=6,
                  bullet=False, level=0):
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    if bullet:
        p.bullet = True
    p.space_before = Pt(space_before)
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def add_footer(slide, text="PULSE · Bago City College · BSIS · 2026"):
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12), Inches(0.3),
                text, size=9, color=MUTED, align=PP_ALIGN.CENTER)


def add_slide_number(slide, num, total=8):
    add_textbox(slide, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.3),
                f"{num}/{total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_section_label(slide, label):
    bar = add_rect(slide, Inches(0.5), Inches(0.35), Inches(0.08), Inches(0.45), TEAL)
    bar.line.fill.background()
    add_textbox(slide, Inches(0.7), Inches(0.3), Inches(4), Inches(0.5),
                label.upper(), size=11, bold=True, color=TEAL)


def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CANVAS)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), PRIMARY)

    # Right accent panel
    add_rect(slide, Inches(8.2), Inches(0.12), Inches(5.13), Inches(7.38), PRIMARY)
    add_rect(slide, Inches(8.5), Inches(1.2), Inches(4.6), Inches(4.8),
             RGBColor(0x0A, 0x3A, 0x62))
    add_textbox(slide, Inches(8.7), Inches(1.5), Inches(4.2), Inches(0.5),
                "OUTBREAK MAP PREVIEW", size=10, bold=True, color=TEAL)
    add_textbox(slide, Inches(8.7), Inches(2.1), Inches(4.2), Inches(3.5),
                "24 Barangay APTAS Choropleth\n\n"
                "● Critical zones (purple)\n"
                "● High-risk zones (red)\n"
                "● Case pins + heatmap\n"
                "● Leaflet + OpenStreetMap",
                size=12, color=WHITE)

    add_textbox(slide, Inches(0.7), Inches(1.0), Inches(7), Inches(0.8),
                "PULSE", size=54, bold=True, color=PRIMARY)
    add_textbox(slide, Inches(0.7), Inches(1.75), Inches(7), Inches(0.5),
                "Public Health Unified Surveillance & Epidemiology",
                size=16, color=TEAL)
    add_textbox(slide, Inches(0.7), Inches(2.35), Inches(7), Inches(0.6),
                "Real-Time Syndromic Surveillance and Risk Forecasting",
                size=20, bold=True, color=DARK)

    tagline_box, tf = add_textbox(
        slide, Inches(0.7), Inches(3.2), Inches(7.2), Inches(1.2),
        "AI-powered disease surveillance that turns barangay health reports into "
        "instant outbreak warnings and geospatial risk maps for faster public health response.",
        size=13, color=DARK)
    tf.paragraphs[0].font.italic = True

    add_rect(slide, Inches(0.7), Inches(4.55), Inches(7), Inches(0.02), TEAL)
    add_textbox(slide, Inches(0.7), Inches(4.75), Inches(7), Inches(0.4),
                "Government Service to Improve:", size=11, bold=True, color=MUTED)
    add_textbox(slide, Inches(0.7), Inches(5.1), Inches(7), Inches(0.5),
                "Bago City Community-Based Disease Surveillance (CDSS)",
                size=15, bold=True, color=PRIMARY)
    add_textbox(slide, Inches(0.7), Inches(5.55), Inches(7), Inches(0.35),
                "City Health Office (CHO) · Negros Occidental", size=12, color=MUTED)

    add_textbox(slide, Inches(0.7), Inches(6.2), Inches(7), Inches(0.6),
                "Banez · Betita · Sarcia  |  Bago City College — BS Information Systems",
                size=11, bold=True, color=DARK)


def slide_02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CANVAS)
    add_section_label(slide, "The Problem")
    add_textbox(slide, Inches(0.7), Inches(0.85), Inches(8), Inches(0.6),
                "Why Bago City's surveillance still struggles today",
                size=28, bold=True, color=PRIMARY)

    problems = [
        ("01", "Inconsistent, Non-Integrated Reporting",
         "Barangay reports arrive at different times, in different formats — slowing consolidation across 24 barangays."),
        ("02", "No Real-Time Analysis or Visualization",
         "Current tools show raw case counts only — no trends, no maps, no risk forecasting for CHO decision-makers."),
        ("03", "No Automated Early Warning System",
         "Potential outbreaks are detected late because there is no system to flag abnormal symptom patterns automatically."),
    ]
    y = 1.65
    for num, title, desc in problems:
        card = add_rect(slide, Inches(0.7), Inches(y), Inches(7.5), Inches(1.15),
                        WHITE, RGBColor(0xFE, 0xCA, 0xCA))
        add_rect(slide, Inches(0.7), Inches(y), Inches(0.12), Inches(1.15), CRITICAL)
        add_textbox(slide, Inches(1.0), Inches(y + 0.1), Inches(0.5), Inches(0.4),
                    num, size=18, bold=True, color=CRITICAL)
        add_textbox(slide, Inches(1.45), Inches(y + 0.08), Inches(6.5), Inches(0.4),
                    title, size=14, bold=True, color=DARK)
        add_textbox(slide, Inches(1.45), Inches(y + 0.48), Inches(6.5), Inches(0.6),
                    desc, size=11, color=MUTED)
        y += 1.3

    # Right column visual placeholder
    add_rect(slide, Inches(8.5), Inches(1.65), Inches(4.2), Inches(2.2),
             RGBColor(0xFE, 0xF2, 0xF2))
    add_textbox(slide, Inches(8.7), Inches(2.0), Inches(3.8), Inches(1.5),
                "CURRENT STATE\n\n📄 CDSS Worksheet (paper)\n⏳ Delayed PIDSR encoding\n📊 Raw counts only",
                size=12, color=CRITICAL, align=PP_ALIGN.CENTER)

    banner = add_rect(slide, Inches(0.5), Inches(5.85), Inches(12.3), Inches(1.15), PRIMARY)
    add_textbox(
        slide, Inches(0.8), Inches(6.0), Inches(11.8), Inches(0.9),
        "Barangay Health Workers and the City Health Office struggle with timely outbreak detection "
        "because disease reporting remains manual, fragmented, and reactive — resulting in delayed "
        "investigations, missed clusters, and weakened community protection across Bago City.",
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_footer(slide)
    add_slide_number(slide, 2)


def slide_03_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CANVAS)
    add_section_label(slide, "Our Solution")
    add_textbox(slide, Inches(0.7), Inches(0.85), Inches(8), Inches(0.5),
                "Solution Overview & Innovation", size=28, bold=True, color=PRIMARY)

    quote_bg = add_rect(slide, Inches(0.7), Inches(1.5), Inches(12), Inches(1.1), RGBColor(0xE0, 0xF2, 0xFE))
    add_textbox(
        slide, Inches(0.95), Inches(1.65), Inches(11.5), Inches(0.9),
        "Our solution uses real-time syndromic data collection, AI-assisted anomaly detection, and "
        "geospatial risk mapping to help Barangay Health Workers and the City Health Office achieve "
        "faster outbreak detection, evidence-based decisions, and coordinated public health response "
        "across all 24 barangays of Bago City.",
        size=13, bold=True, color=PRIMARY)

    steps = [
        ("CAPTURE", "BHW digital\nCDSS submission", TEAL),
        ("ANALYZE", "ML screening +\nAPTAS scoring", PRIMARY),
        ("ACT", "Alerts + map +\nCHO response", ACCENT),
    ]
    x = 1.2
    for label, desc, color in steps:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.9), Inches(1.2), Inches(1.2))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        add_textbox(slide, Inches(x), Inches(3.2), Inches(1.2), Inches(0.5),
                    label, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x - 0.2), Inches(4.2), Inches(1.6), Inches(0.8),
                    desc, size=11, color=DARK, align=PP_ALIGN.CENTER)
        if x < 8:
            add_textbox(slide, Inches(x + 1.3), Inches(3.35), Inches(0.5), Inches(0.4),
                        "→", size=22, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
        x += 3.8

    add_textbox(slide, Inches(0.7), Inches(5.1), Inches(6), Inches(0.4),
                "Innovation: Unlike PIDSR & DHIS2", size=14, bold=True, color=DARK)
    innovations = [
        "✓  Instant ML screening on every case submission",
        "✓  Multi-variate APTAS risk scoring (A+T+E+S)",
        "✓  Live 24-barangay choropleth outbreak map",
        "✓  Role-based dashboards for BHWs & CHO staff",
    ]
    y = 5.5
    for item in innovations:
        add_textbox(slide, Inches(0.7), Inches(y), Inches(6.2), Inches(0.3),
                    item, size=11, color=DARK)
        y += 0.32

    add_rect(slide, Inches(7.3), Inches(5.0), Inches(5.5), Inches(1.7), PRIMARY)
    add_textbox(slide, Inches(7.5), Inches(5.2), Inches(5.1), Inches(1.3),
                "From paper worksheets\n→ digital intelligence\nin one platform.",
                size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_footer(slide)
    add_slide_number(slide, 3)


def slide_04_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CANVAS)
    add_section_label(slide, "Core Capabilities")
    add_textbox(slide, Inches(0.7), Inches(0.85), Inches(8), Inches(0.5),
                "4 Features That Deliver Impact", size=28, bold=True, color=PRIMARY)

    features = [
        ("Digital Syndromic\nCase Submission",
         "BHWs encode CDSS Worksheet data — demographics, 39 PIDSR symptom codes, GPS pin — in one batch form",
         "Cuts double-encoding and eliminates paper-to-PIDSR lag",
         TEAL),
        ("Real-Time Anomaly\nDetection",
         "Isolation Forest screens every submission; Random Forest classifies syndrome type on intake",
         "Flags probable outbreaks at submission — not days later",
         PRIMARY),
        ("APTAS Multi-Variate\nAlerting",
         "Combines Anomaly + Temporal + Environmental + Spatial scores into 0–100 risk index",
         "CHO gets actionable early warnings with circuit-breaker logic",
         WARNING),
        ("Geospatial Hotspot\nMapping",
         "Leaflet + OSM choropleth across 24 barangays — pins, heatmap, risk legend",
         "Officers see exactly where to deploy field teams",
         ACCENT),
    ]

    x = 0.55
    for name, does, benefit, color in features:
        card = add_rect(slide, Inches(x), Inches(1.55), Inches(2.95), Inches(5.0), WHITE, RGBColor(0xE2, 0xE8, 0xF0))
        add_rect(slide, Inches(x), Inches(1.55), Inches(2.95), Inches(0.55), color)
        add_textbox(slide, Inches(x + 0.15), Inches(1.62), Inches(2.65), Inches(0.7),
                    name, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.15), Inches(2.25), Inches(2.65), Inches(0.35),
                    "WHAT IT DOES", size=8, bold=True, color=MUTED)
        add_textbox(slide, Inches(x + 0.15), Inches(2.55), Inches(2.65), Inches(1.5),
                    does, size=10, color=DARK)
        add_textbox(slide, Inches(x + 0.15), Inches(4.2), Inches(2.65), Inches(0.35),
                    "BENEFIT", size=8, bold=True, color=ACCENT)
        add_textbox(slide, Inches(x + 0.15), Inches(4.5), Inches(2.65), Inches(1.8),
                    benefit, size=11, bold=True, color=DARK)
        x += 3.15

    add_footer(slide)
    add_slide_number(slide, 4)


def slide_05_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CANVAS)
    add_section_label(slide, "Architecture & Data Flow")
    add_textbox(slide, Inches(0.7), Inches(0.85), Inches(10), Inches(0.5),
                "From Paper Records to Digital Outbreak Maps", size=26, bold=True, color=PRIMARY)

    layers = [
        ("Presentation", "Web dashboards · Landing page · Outbreak map (Django + Leaflet)", PRIMARY),
        ("Application", "Role-based auth — BHW · Encoder · Admin · Surveillance Officer", TEAL),
        ("Analytics", "Isolation Forest + Random Forest · APTAS engine (A+T+E+S)", WARNING),
        ("Integration", "Open-Meteo weather API · 24-barangay GeoJSON boundaries", ACCENT),
        ("Data", "MySQL (pulse_db) · Cloud-hosted · Secure backup", DARK),
    ]
    y = 1.45
    for name, desc, color in layers:
        add_rect(slide, Inches(0.7), Inches(y), Inches(7.5), Inches(0.62), color)
        add_textbox(slide, Inches(0.9), Inches(y + 0.05), Inches(1.5), Inches(0.5),
                    name, size=11, bold=True, color=WHITE)
        add_textbox(slide, Inches(2.5), Inches(y + 0.1), Inches(5.5), Inches(0.45),
                    desc, size=10, color=WHITE)
        y += 0.7

    add_textbox(slide, Inches(0.7), Inches(5.0), Inches(7.5), Inches(0.35),
                "DATA FLOW PIPELINE", size=10, bold=True, color=MUTED)

    flow = (
        "BHW Field Data → Digital Submission → ML Screening → APTAS Scoring → "
        "Dashboard Alerts → Choropleth Map → CHO Investigation"
    )
    flow_bg = add_rect(slide, Inches(0.7), Inches(5.35), Inches(12), Inches(0.75), RGBColor(0xE0, 0xF2, 0xFE))
    add_textbox(slide, Inches(0.9), Inches(5.5), Inches(11.6), Inches(0.5),
                flow, size=11, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(8.5), Inches(1.45), Inches(4.2), Inches(3.3), WHITE, RGBColor(0xE2, 0xE8, 0xF0))
    add_textbox(slide, Inches(8.7), Inches(1.6), Inches(3.8), Inches(0.35),
                "KEY DATABASE ENTITIES", size=10, bold=True, color=PRIMARY)
    entities = "admins · bhw · health_reports · patients · barangays · environment_data · ml_predictions · risk_analysis · alerts"
    add_textbox(slide, Inches(8.7), Inches(2.0), Inches(3.8), Inches(2.5), entities, size=10, color=DARK)

    add_footer(slide)
    add_slide_number(slide, 5)


def slide_06_prototype(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CANVAS)
    add_section_label(slide, "Live Prototype")
    add_textbox(slide, Inches(0.7), Inches(0.85), Inches(10), Inches(0.5),
                "Prototype Walkthrough — Deployment Ready", size=26, bold=True, color=PRIMARY)

    screens = [
        ("① Landing Page", "Hero: AI Assistive Disease Surveillance\nFeatures · Staff sign-in · CHO branding", PRIMARY),
        ("② Login / Auth", "Role-based access · Account lockout\nFirst-login password setup", TEAL),
        ("③ Admin Dashboard", "Morbidity trends · APTAS alert panel\nPending validations · Weather widget", WARNING),
        ("④ BHW Dashboard", "Assigned barangay locked\nRecent submissions · APTAS cards", ACCENT),
        ("⑤ Outbreak Map", "24-barangay choropleth\nCritical & High zones · Case pins · Heatmap", CRITICAL),
        ("⑥ Risk Analysis", "Syndrome filters · Risk timeline\nInvestigation status · Incidents", DARK),
    ]

    positions = [(0.55, 1.5), (4.55, 1.5), (8.55, 1.5), (0.55, 4.2), (4.55, 4.2), (8.55, 4.2)]
    for (x, y), (title, desc, color) in zip(positions, screens):
        card = add_rect(slide, Inches(x), Inches(y), Inches(3.7), Inches(2.4), WHITE, RGBColor(0xE2, 0xE8, 0xF0))
        add_rect(slide, Inches(x), Inches(y), Inches(3.7), Inches(0.45), color)
        add_textbox(slide, Inches(x + 0.15), Inches(y + 0.08), Inches(3.4), Inches(0.35),
                    title, size=11, bold=True, color=WHITE)
        # Screenshot placeholder
        add_rect(slide, Inches(x + 0.15), Inches(y + 0.55), Inches(3.4), Inches(1.2),
                 RGBColor(0xF1, 0xF5, 0xF9))
        add_textbox(slide, Inches(x + 0.15), Inches(y + 0.85), Inches(3.4), Inches(0.5),
                    "[ Insert screenshot ]", size=9, color=MUTED, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.15), Inches(y + 1.85), Inches(3.4), Inches(0.5),
                    desc, size=9, color=DARK)

    add_textbox(slide, Inches(0.7), Inches(6.75), Inches(12), Inches(0.3),
                "✅ Core modules functional  ·  ISO/IEC 25010 evaluated  ·  Cloud deployment configured",
                size=10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    add_footer(slide)
    add_slide_number(slide, 6)


def slide_07_business(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CANVAS)
    add_section_label(slide, "Sustainability & Impact")
    add_textbox(slide, Inches(0.7), Inches(0.85), Inches(10), Inches(0.5),
                "Business Model · Traction · Funds · SDGs", size=26, bold=True, color=PRIMARY)

    # Business model box
    add_rect(slide, Inches(0.7), Inches(1.45), Inches(5.8), Inches(2.0), WHITE, RGBColor(0xE2, 0xE8, 0xF0))
    add_textbox(slide, Inches(0.9), Inches(1.55), Inches(5.4), Inches(0.35),
                "G2G / Localized SaaS Model", size=13, bold=True, color=PRIMARY)
    model_items = [
        "• Annual LGU license for Bago City CHO + barangay stations",
        "• Replicable across Negros Occidental LGUs",
        "• Complements PIDSR — does not replace national reporting",
        "• Cloud-hosted: Vercel/AWS · MySQL · Multi-tenant ready",
    ]
    y = 1.95
    for item in model_items:
        add_textbox(slide, Inches(0.9), Inches(y), Inches(5.4), Inches(0.3), item, size=10, color=DARK)
        y += 0.35

    # Traction box
    add_rect(slide, Inches(6.8), Inches(1.45), Inches(5.8), Inches(2.0), RGBColor(0xEC, 0xFD, 0xF5))
    add_textbox(slide, Inches(7.0), Inches(1.55), Inches(5.4), Inches(0.35),
                "Traction & Validation", size=13, bold=True, color=ACCENT)
    traction = [
        "✅ Fully developed system (auth, ML, APTAS, map)",
        "✅ CDSS Unit stakeholder-aligned workflows",
        "✅ ISO/IEC 25010:2011 quality evaluation",
        "✅ 24 barangays mapped · Deployment-ready",
        "✅ Agile/Scrum with CHO feedback loops",
    ]
    y = 1.95
    for item in traction:
        add_textbox(slide, Inches(7.0), Inches(y), Inches(5.4), Inches(0.3), item, size=10, color=DARK)
        y += 0.35

    # Funds table
    add_textbox(slide, Inches(0.7), Inches(3.6), Inches(5), Inches(0.35),
                "Proposed Use of Funds", size=13, bold=True, color=PRIMARY)
    funds = [
        ("Cloud Hosting & Security", "25%", "Uptime, SSL, encrypted backups, Data Privacy Act"),
        ("ML Model Fine-Tuning", "25%", "Retrain on live Bago City morbidity data"),
        ("BHW & CHO Training", "20%", "Field orientation, manuals, barangay rollout"),
        ("LGU & DOH Integration", "15%", "PIDSR bridge, RESU connectivity, SMS alerts"),
        ("UI/UX & Mobile Access", "15%", "Responsive design for low-bandwidth stations"),
    ]
    y = 4.0
    for name, pct, purpose in funds:
        add_rect(slide, Inches(0.7), Inches(y), Inches(6.0), Inches(0.42), WHITE, RGBColor(0xE2, 0xE8, 0xF0))
        add_textbox(slide, Inches(0.85), Inches(y + 0.06), Inches(2.8), Inches(0.3),
                    name, size=10, bold=True, color=DARK)
        add_textbox(slide, Inches(3.7), Inches(y + 0.06), Inches(0.6), Inches(0.3),
                    pct, size=11, bold=True, color=TEAL)
        add_textbox(slide, Inches(4.4), Inches(y + 0.06), Inches(2.2), Inches(0.3),
                    purpose, size=9, color=MUTED)
        y += 0.48

    # SDGs
    add_rect(slide, Inches(7.0), Inches(3.6), Inches(5.6), Inches(2.9), PRIMARY)
    add_textbox(slide, Inches(7.2), Inches(3.75), Inches(5.2), Inches(0.35),
                "Target Sustainable Development Goals", size=12, bold=True, color=WHITE)
    sdgs = [
        ("SDG 3", "Good Health & Well-Being", "Early outbreak detection protects communities"),
        ("SDG 11", "Sustainable Cities", "Resilient urban health surveillance for Bago City"),
        ("SDG 16", "Strong Institutions", "Strengthens LGU health governance & transparency"),
    ]
    y = 4.2
    for num, title, desc in sdgs:
        add_textbox(slide, Inches(7.2), Inches(y), Inches(1.0), Inches(0.3),
                    num, size=12, bold=True, color=TEAL)
        add_textbox(slide, Inches(8.2), Inches(y), Inches(4.2), Inches(0.55),
                    f"{title}\n{desc}", size=10, color=WHITE)
        y += 0.75

    add_footer(slide)
    add_slide_number(slide, 7)


def slide_08_team(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CANVAS)
    add_section_label(slide, "The Team")
    add_textbox(slide, Inches(0.7), Inches(0.85), Inches(8), Inches(0.5),
                "Development Team & Closing", size=28, bold=True, color=PRIMARY)

    team = [
        ("Hernani J. Banez", "Backend · Django · Database design"),
        ("Jaykirt M. Betita", "ML pipeline · APTAS · Risk analytics"),
        ("Earl Adrian M. Sarcia", "Frontend · GIS mapping · UI/UX"),
    ]
    x = 0.7
    for name, role in team:
        card = add_rect(slide, Inches(x), Inches(1.5), Inches(3.7), Inches(2.2), WHITE, RGBColor(0xE2, 0xE8, 0xF0))
        avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.35), Inches(1.7), Inches(1.0), Inches(1.0))
        avatar.fill.solid()
        avatar.fill.fore_color.rgb = PRIMARY
        avatar.line.fill.background()
        add_textbox(slide, Inches(x + 0.15), Inches(2.85), Inches(3.4), Inches(0.35),
                    name, size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.15), Inches(3.2), Inches(3.4), Inches(0.4),
                    role, size=10, color=MUTED, align=PP_ALIGN.CENTER)
        x += 4.0

    add_rect(slide, Inches(0.7), Inches(4.0), Inches(5.5), Inches(1.5), RGBColor(0xE0, 0xF2, 0xFE))
    add_textbox(slide, Inches(0.9), Inches(4.15), Inches(5.1), Inches(1.2),
                "Bago City College\nBachelor of Science in Information Systems\n"
                "Project Coordinator: Anthony S. Malabanan, MIT, MAT-Math\n"
                "Partner: Bago City Health Office — CDSS Unit",
                size=11, color=PRIMARY)

    add_rect(slide, Inches(6.5), Inches(4.0), Inches(5.8), Inches(1.5), RGBColor(0xF1, 0xF5, 0xF9))
    add_textbox(slide, Inches(6.7), Inches(4.2), Inches(5.4), Inches(1.1),
                "LIVE DEMO\n\n[ Insert QR code / URL ]\npulse-ai.vercel.app",
                size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    banner = add_rect(slide, Inches(0.5), Inches(5.75), Inches(12.3), Inches(1.25), PRIMARY)
    add_textbox(
        slide, Inches(0.8), Inches(5.95), Inches(11.8), Inches(0.9),
        "PULSE transforms delayed, paper-based barangay disease reporting into real-time, "
        "map-driven outbreak intelligence — giving Bago City the early warning system its communities deserve.",
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.7), Inches(7.0), Inches(12), Inches(0.35),
                "Thank you. We welcome your questions.", size=12, bold=True,
                color=TEAL, align=PP_ALIGN.CENTER)

    add_slide_number(slide, 8)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_features(prs)
    slide_05_architecture(prs)
    slide_06_prototype(prs)
    slide_07_business(prs)
    slide_08_team(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    build()
